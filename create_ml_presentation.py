"""
create_ml_presentation.py
=========================
Generates the comprehensive ML/ERG classification presentation as a PPTX file.
Covers the full pipeline: motivation → chirp → NI → dual → CNN interpretability.

Run from /Users/leo/retina/machine_learning/:
    python create_ml_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx.oxml.ns as qn
from lxml import etree
from copy import deepcopy

# ── Paths ─────────────────────────────────────────────────────────────────────
ML_ROOT   = os.path.dirname(os.path.abspath(__file__))
CHIRP_FIG = os.path.join(ML_ROOT, 'chirp_analysis', 'results', 'figures')
NI_FIG    = os.path.join(ML_ROOT, 'natural_image_analysis', 'results', 'figures')
MS_ROOT   = os.path.abspath(os.path.join(ML_ROOT, '..', 'multi_stimuli'))
MS_FIG    = os.path.join(MS_ROOT, 'results')
OUT_PATH  = os.path.join(ML_ROOT, 'ML_ERG_Full_Presentation.pptx')

def cf(name):   return os.path.join(CHIRP_FIG, name)
def nf(name):   return os.path.join(NI_FIG, name)
def mf(*parts): return os.path.join(MS_FIG, *parts)

# ── Design constants ───────────────────────────────────────────────────────────
W, H = Inches(13.333), Inches(7.5)   # 16:9 widescreen

# Color palette
C_NAVY    = RGBColor(0x1E, 0x3A, 0x5F)   # dark navy – headers
C_BLUE    = RGBColor(0x27, 0x6F, 0xBF)   # medium blue
C_GOLD    = RGBColor(0xE8, 0xA8, 0x38)   # amber/gold – highlights
C_RED     = RGBColor(0xD6, 0x28, 0x28)   # 5xFAD red
C_WTBLUE  = RGBColor(0x21, 0x96, 0xF3)   # WT blue
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)   # success green
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_OFFWT   = RGBColor(0xF8, 0xF9, 0xFA)   # near-white bg
C_DARKGR  = RGBColor(0x2C, 0x3E, 0x50)   # dark grey text
C_MIDGR   = RGBColor(0x5D, 0x6D, 0x7E)   # medium grey
C_LIGHTGR = RGBColor(0xEC, 0xF0, 0xF1)   # light grey for boxes
C_DIVBG   = RGBColor(0x12, 0x24, 0x3E)   # section divider bg (deeper navy)

FONT_TITLE  = 'Calibri'
FONT_BODY   = 'Calibri'

# ── Master helpers ─────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE.RECTANGLE
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_name=FONT_BODY, size=Pt(14), bold=False, italic=False,
             color=C_DARKGR, align=PP_ALIGN.LEFT, wrap=True,
             fill=None, line=None):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = size
    run.font.bold  = bold
    run.font.italic= italic
    run.font.color.rgb = color
    if fill:
        txb.fill.solid()
        txb.fill.fore_color.rgb = fill
    if line:
        txb.line.color.rgb = line
    return txb

def add_image(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        print(f"  [WARN] Missing: {path}")
        add_rect(slide, x, y, w, h or Inches(3),
                 fill=RGBColor(0xDD,0xDD,0xDD))
        add_text(slide, f"[{os.path.basename(path)}]",
                 x, y+(h or Inches(3))/2 - Inches(0.2), w, Inches(0.4),
                 size=Pt(10), color=C_MIDGR, align=PP_ALIGN.CENTER)
        return
    if h:
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        slide.shapes.add_picture(path, x, y, w)

def header_bar(slide, title, subtitle=None):
    """Dark navy header bar at top of content slides."""
    bar_h = Inches(1.05)
    add_rect(slide, 0, 0, W, bar_h, fill=C_NAVY)
    # Accent line
    add_rect(slide, 0, bar_h - Inches(0.05), W, Inches(0.05), fill=C_GOLD)
    # Title text
    add_text(slide, title,
             Inches(0.35), Inches(0.08), W - Inches(0.7), Inches(0.6),
             size=Pt(28), bold=True, color=C_WHITE,
             font_name=FONT_TITLE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.62), W - Inches(0.7), Inches(0.38),
                 size=Pt(14), bold=False, color=RGBColor(0xB0,0xC8,0xE8),
                 font_name=FONT_BODY, align=PP_ALIGN.LEFT)
    return bar_h

def slide_number_tag(slide, n):
    """Small slide number bottom-right."""
    add_text(slide, str(n),
             W - Inches(0.5), H - Inches(0.32), Inches(0.4), Inches(0.28),
             size=Pt(9), color=C_MIDGR, align=PP_ALIGN.RIGHT)

def bullet_box(slide, items, x, y, w, h,
               title=None, title_size=Pt(15), body_size=Pt(13),
               bg=None, border=None):
    """Creates a box with optional title and bullet list."""
    if bg:
        add_rect(slide, x, y, w, h, fill=bg, line=border, line_w=Pt(1))
    if title:
        add_text(slide, title, x+Inches(0.1), y+Inches(0.06),
                 w-Inches(0.2), Inches(0.38),
                 size=title_size, bold=True, color=C_NAVY)
        y0 = y + Inches(0.42)
    else:
        y0 = y + Inches(0.1)
    bh = (h - (y0-y) - Inches(0.05)) / max(len(items),1)
    for item in items:
        add_text(slide, f"• {item}", x+Inches(0.15), y0, w-Inches(0.3), bh,
                 size=body_size, color=C_DARKGR)
        y0 += bh

def metric_badge(slide, label, value, x, y, w=Inches(1.8), h=Inches(0.95),
                 bg=C_NAVY, val_color=C_GOLD):
    """Metric badge: colored box with label + large value."""
    add_rect(slide, x, y, w, h, fill=bg)
    add_text(slide, label, x, y+Inches(0.04), w, Inches(0.35),
             size=Pt(11), color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, value, x, y+Inches(0.38), w, Inches(0.52),
             size=Pt(22), bold=True, color=val_color, align=PP_ALIGN.CENTER)

# ── Slide 1 – Title ────────────────────────────────────────────────────────────
def slide_title(prs, n):
    s = blank_slide(prs)
    # Full background
    add_rect(s, 0, 0, W, H, fill=C_DIVBG)
    # Decorative stripe
    add_rect(s, 0, H*0.58, W, Inches(0.08), fill=C_GOLD)
    # Gold accent bar left
    add_rect(s, 0, 0, Inches(0.15), H, fill=C_GOLD)
    # Title
    add_text(s, "Machine Learning Classification of Alzheimer's Disease\nfrom Electroretinogram Responses",
             Inches(0.5), Inches(1.2), W-Inches(1.0), Inches(2.0),
             size=Pt(36), bold=True, color=C_WHITE, font_name=FONT_TITLE,
             align=PP_ALIGN.LEFT)
    # Subtitle
    add_text(s,
             "CNN, Handcrafted ML, Complexity Features & Interpretability Analysis\n"
             "Binary WT vs 5xFAD Classification from Chirp & Natural Image ERG",
             Inches(0.5), Inches(3.3), W-Inches(1.0), Inches(1.1),
             size=Pt(18), color=RGBColor(0xB0,0xC8,0xE8), align=PP_ALIGN.LEFT)
    # Stimuli labels
    add_text(s, "Chirp  ·  Natural Image  ·  Dual Fusion",
             Inches(0.5), Inches(4.5), W-Inches(1.0), Inches(0.5),
             size=Pt(15), color=C_GOLD, align=PP_ALIGN.LEFT)
    # Author/date
    add_text(s, "Leo Medina  ·  2026",
             Inches(0.5), H-Inches(0.7), Inches(4), Inches(0.45),
             size=Pt(13), color=C_MIDGR, align=PP_ALIGN.LEFT)
    slide_number_tag(s, n)

# ── Slide 2 – Motivation ───────────────────────────────────────────────────────
def slide_motivation(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Motivation",
                       "Why use machine learning on retinal electrophysiology?")
    Y = bar_h + Inches(0.2)

    # Left column – biological context
    add_rect(s, Inches(0.25), Y, Inches(6.1), H-Y-Inches(0.35),
             fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(1))
    add_text(s, "Biological Context", Inches(0.4), Y+Inches(0.1),
             Inches(5.8), Inches(0.4), size=Pt(16), bold=True, color=C_NAVY)
    ctx = [
        "Alzheimer's disease (AD) causes progressive neurodegeneration,\n"
        "  including in the retina — the accessible 'window to the brain'",
        "5xFAD transgenic mice overexpress human amyloid precursors,\n"
        "  developing AD-like pathology by 2 months of age",
        "Electroretinogram (ERG) records mass retinal electrical responses\n"
        "  non-invasively — potentially a cheap, fast AD biomarker",
        "Previous study (Araya-Arriagada et al., 2022) showed ERG differences\n"
        "  are detectable with traditional statistics in adult mice only",
        "Goal: Can ML detect subtle patterns even in young animals?\n"
        "  And what do those patterns tell us biologically?",
    ]
    yy = Y + Inches(0.6)
    for c in ctx:
        add_text(s, f"• {c}", Inches(0.4), yy, Inches(5.8), Inches(0.64),
                 size=Pt(12.5), color=C_DARKGR)
        yy += Inches(0.64)

    # Right column – approach
    add_rect(s, Inches(6.6), Y, Inches(6.48), H-Y-Inches(0.35),
             fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(1))
    add_text(s, "Study Approach", Inches(6.75), Y+Inches(0.1),
             Inches(6.1), Inches(0.4), size=Pt(16), bold=True, color=C_NAVY)
    appr = [
        ("1. CNN End-to-End",     "Train deep 1D-CNN directly on raw ERG traces"),
        ("2. Handcrafted ML",     "Physiologically motivated features + classical classifiers"),
        ("3. Complexity Features","Sample entropy / MSE curves capturing nonlinear dynamics"),
        ("4. Interpretability",   "Understand WHAT the CNN learns → new biological insight"),
        ("5. Multi-stimulus",     "Combine chirp + natural image ERG for richer representation"),
    ]
    yy = Y + Inches(0.6)
    for (lbl, desc) in appr:
        add_rect(s, Inches(6.7), yy, Inches(0.12), Inches(0.35), fill=C_GOLD)
        add_text(s, lbl, Inches(6.9), yy, Inches(2.3), Inches(0.35),
                 size=Pt(12.5), bold=True, color=C_NAVY)
        add_text(s, desc, Inches(9.2), yy, Inches(3.7), Inches(0.35),
                 size=Pt(12), color=C_DARKGR)
        yy += Inches(0.6)
    slide_number_tag(s, n)

# ── Slide 3 – Experimental Setup ───────────────────────────────────────────────
def slide_experiment(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Experimental Setup",
                       "In-vitro micro-ERG recording from isolated retinal preparations")
    add_image(s, mf('representative_traces.png'),
              Inches(0.25), bar_h+Inches(0.15),
              Inches(8.5), H-bar_h-Inches(0.5))
    # Info boxes right
    X = Inches(9.0); Y = bar_h+Inches(0.15)
    bh = (H-Y-Inches(0.4)) / 4
    for (lbl, val) in [
        ("Preparations", "N = 46"),
        ("Animals", "38 independent"),
        ("Groups", "4 (WT/5xFAD × Young/Adult)"),
        ("Recording", "In-vitro micro-array ERG"),
    ]:
        add_rect(s, X, Y, Inches(4.08), bh-Inches(0.06), fill=C_NAVY)
        add_text(s, lbl, X+Inches(0.1), Y+Inches(0.04), Inches(3.8), Inches(0.32),
                 size=Pt(12), color=RGBColor(0xB0,0xC8,0xE8))
        add_text(s, val, X+Inches(0.1), Y+Inches(0.36), Inches(3.8), Inches(0.44),
                 size=Pt(15), bold=True, color=C_GOLD)
        Y += bh
    slide_number_tag(s, n)

# ── Slide 4 – Demographics ─────────────────────────────────────────────────────
def slide_demographics(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Demographics & Study Groups")
    add_image(s, mf('demographics_distribution.png'),
              Inches(0.25), bar_h+Inches(0.1),
              Inches(7.0), H-bar_h-Inches(0.4))

    # Table of group stats
    X = Inches(7.5); Y = bar_h+Inches(0.2)
    add_text(s, "Group Summary", X, Y, Inches(5.6), Inches(0.38),
             size=Pt(15), bold=True, color=C_NAVY)
    Y += Inches(0.42)
    headers = ["Group","N","Age (months)","Sex"]
    rows = [
        ["WT adult",   "10", "6.9 ± 0.5", "6F/4M"],
        ["WT young",   "13", "2.6 ± 0.6", "8F/5M"],
        ["5xFAD adult","14", "7.0 ± 0.6", "8F/6M"],
        ["5xFAD young"," 9", "2.5 ± 0.3", "5F/4M"],
        ["TOTAL",      "46", "–",          "–"],
    ]
    col_w = [Inches(1.7), Inches(0.6), Inches(1.5), Inches(1.0)]
    rh = Inches(0.42)
    # Header row
    cx = X
    for (h, cw) in zip(headers, col_w):
        add_rect(s, cx, Y, cw-Inches(0.03), rh, fill=C_NAVY)
        add_text(s, h, cx+Inches(0.05), Y+Inches(0.05), cw-Inches(0.1), rh-Inches(0.1),
                 size=Pt(11), bold=True, color=C_WHITE)
        cx += cw
    Y += rh
    for i, row in enumerate(rows):
        bg = C_OFFWT if i%2==0 else C_WHITE
        if row[0] == "TOTAL": bg = C_LIGHTGR
        cx = X
        for (cell, cw) in zip(row, col_w):
            add_rect(s, cx, Y, cw-Inches(0.03), rh-Inches(0.02),
                     fill=bg, line=C_LIGHTGR, line_w=Pt(0.5))
            add_text(s, cell, cx+Inches(0.05), Y+Inches(0.06), cw-Inches(0.1), rh-Inches(0.12),
                     size=Pt(11), bold=(row[0]=="TOTAL"), color=C_DARKGR)
            cx += cw
        Y += rh

    add_text(s, "• 8 animals contributed two preparations (−t1/−t2 pairs)\n"
               "• Subject-disjoint cross-validation used throughout",
             X, Y+Inches(0.1), Inches(5.6), Inches(0.7),
             size=Pt(11), color=C_MIDGR)
    slide_number_tag(s, n)

# ── Slide 5 – Stimuli ──────────────────────────────────────────────────────────
def slide_stimuli(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "ERG Stimuli",
                       "Two complementary stimuli probe different aspects of retinal function")
    add_image(s, mf('chirp_traces.png'),
              Inches(0.25), bar_h+Inches(0.1),
              Inches(6.5), H-bar_h-Inches(1.5))
    add_image(s, mf('ni_traces.png'),
              Inches(6.9), bar_h+Inches(0.1),
              Inches(6.2), H-bar_h-Inches(1.5))

    # Labels below
    BY = H - Inches(1.3)
    add_rect(s, Inches(0.25), BY, Inches(6.5), Inches(0.9), fill=C_NAVY)
    add_text(s, "CHIRP stimulus (35 s, 250 Hz sampling)",
             Inches(0.4), BY+Inches(0.05), Inches(6.2), Inches(0.38),
             size=Pt(13), bold=True, color=C_GOLD)
    add_text(s, "Flash → Frequency ramp (0→8 Hz) → Amplitude ramp (1 Hz)",
             Inches(0.4), BY+Inches(0.42), Inches(6.2), Inches(0.38),
             size=Pt(11), color=C_WHITE)

    add_rect(s, Inches(6.9), BY, Inches(6.2), Inches(0.9), fill=C_NAVY)
    add_text(s, "NATURAL IMAGE stimulus (10 reps, 250 Hz sampling)",
             Inches(7.05), BY+Inches(0.05), Inches(5.9), Inches(0.38),
             size=Pt(13), bold=True, color=C_GOLD)
    add_text(s, "10 × 2500 samples — spatial / luminance pattern flashes",
             Inches(7.05), BY+Inches(0.42), Inches(5.9), Inches(0.38),
             size=Pt(11), color=C_WHITE)
    slide_number_tag(s, n)

# ── Slide 6 – Chirp Segments ───────────────────────────────────────────────────
def slide_chirp_segments(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Chirp Segments",
                       "The 35-second chirp is divided into functionally distinct epochs")
    add_image(s, cf('17_amplitude_inspection.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(8.8), H-bar_h-Inches(0.35))

    X = Inches(9.2); Y = bar_h+Inches(0.15)
    segs = [
        ("full",           "0 – 35 s",   "8750 pts", "All segments combined"),
        ("flash",          "0 – 7.5 s",  "1875 pts", "Initial flash transient"),
        ("frequency",      "7.5 – 24 s", "4125 pts", "Frequency chirp (0→8 Hz)"),
        ("amplitude",      "24 – 35 s",  "2750 pts", "★  Amplitude ramp (1 Hz)\n   ← Best for CNN"),
        ("amplitude_norm", "24 – 35 s",  "2750 pts", "Amplitude normalized\n   by RMS"),
    ]
    colors = [C_BLUE, C_BLUE, C_BLUE, C_GOLD, C_BLUE]
    bh = (H-Y-Inches(0.35)) / len(segs)
    for (nm, times, pts, desc), col in zip(segs, colors):
        add_rect(s, X, Y, Inches(0.12), bh-Inches(0.06), fill=col)
        add_text(s, nm, X+Inches(0.18), Y+Inches(0.01),
                 Inches(1.6), Inches(0.32), size=Pt(12), bold=True, color=C_NAVY)
        add_text(s, times, X+Inches(1.8), Y+Inches(0.01),
                 Inches(1.0), Inches(0.32), size=Pt(11), color=C_MIDGR)
        add_text(s, pts, X+Inches(2.8), Y+Inches(0.01),
                 Inches(0.9), Inches(0.32), size=Pt(11), color=C_MIDGR)
        add_text(s, desc, X+Inches(0.18), Y+Inches(0.34),
                 Inches(3.95), bh-Inches(0.38), size=Pt(10.5), color=C_DARKGR)
        Y += bh
    slide_number_tag(s, n)

# ── Slide 7 – Handcrafted Features ────────────────────────────────────────────
def slide_hc_features(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Hand-Crafted ERG Features",
                       "13 physiologically motivated features derived from flash and chirp responses")
    add_image(s, cf('10_chirp_hc_stats_distribution.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(9.3), H-bar_h-Inches(0.35))

    X = Inches(9.6); Y = bar_h+Inches(0.15)
    features = [
        ("Flash Features (4)", [
            "Flash_Peak_Max — maximum amplitude",
            "Flash_Peak_Min — minimum amplitude",
            "Flash_Peak_P2P — peak-to-peak",
            "Flash_RMS — root-mean-square"]),
        ("Frequency Segment (2)", [
            "ChirpFreq_RMS",
            "ChirpFreq_Std"]),
        ("Amplitude Segment (3)", [
            "ChirpAmp_RMS",
            "ChirpAmp_Max",
            "ChirpAmp_P2P"]),
        ("Power Spectrum (4)", [
            "Power_Total",
            "Power_Low (<2 Hz)",
            "Power_Mid (2–8 Hz)",
            "Power_High (>8 Hz)"]),
    ]
    bh = (H-Y-Inches(0.35)) / len(features)
    for grp, feats in features:
        add_text(s, grp, X, Y+Inches(0.02), Inches(3.7), Inches(0.33),
                 size=Pt(12), bold=True, color=C_NAVY)
        yf = Y + Inches(0.35)
        fh = (bh - Inches(0.38)) / len(feats)
        for f in feats:
            add_text(s, f"  · {f}", X+Inches(0.1), yf, Inches(3.6), fh,
                     size=Pt(10.5), color=C_DARKGR)
            yf += fh
        Y += bh
    slide_number_tag(s, n)

# ── Slide 8 – Complexity Features ─────────────────────────────────────────────
def slide_complexity(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Complexity Features — MSE Analysis",
                       "Multi-scale sample entropy captures nonlinear dynamical structure")
    add_image(s, cf('11_chirp_comp_stats_distribution.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(9.5), H-bar_h-Inches(0.35))

    X = Inches(9.8); Y = bar_h+Inches(0.2)
    add_text(s, "Multi-Scale Entropy (MSE)", X, Y, Inches(3.4), Inches(0.38),
             size=Pt(14), bold=True, color=C_NAVY)
    Y += Inches(0.42)
    items = [
        "Sample entropy SampEn(r, m, τ) at\n  scale factors τ = 1…20",
        "r = 0.2·SD(signal)\nm = 2 (embedding dimension)",
        "Captures complexity at slow\n  vs fast timescales",
        "Features: SampEn values binned\n  into Low / Mid / High scale bands",
        "Fitted MSE curve parameters:\n  slope, plateau, area-under-curve",
    ]
    for item in items:
        add_text(s, f"• {item}", X, Y, Inches(3.4), Inches(0.62),
                 size=Pt(11.5), color=C_DARKGR)
        Y += Inches(0.62)

    # MSE figure for separate panels
    add_image(s, mf('Fig2_MSE_Curves', 'Figure2_MSE_Curves_1D.png'),
              X, Y+Inches(0.05), Inches(3.4))
    slide_number_tag(s, n)

# ── Section divider ────────────────────────────────────────────────────────────
def slide_section(prs, n, part, title, subtitle):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill=C_DIVBG)
    add_rect(s, 0, H*0.44, W, Inches(0.07), fill=C_GOLD)
    add_rect(s, 0, 0, Inches(0.15), H, fill=C_GOLD)
    add_text(s, part, Inches(0.5), Inches(1.4), W-Inches(1), Inches(0.55),
             size=Pt(22), color=C_GOLD, bold=False, align=PP_ALIGN.LEFT)
    add_text(s, title, Inches(0.5), Inches(2.0), W-Inches(1), Inches(1.4),
             size=Pt(40), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT,
             font_name=FONT_TITLE)
    add_text(s, subtitle, Inches(0.5), Inches(3.55), W-Inches(1), Inches(0.7),
             size=Pt(17), color=RGBColor(0xB0,0xC8,0xE8), align=PP_ALIGN.LEFT)
    slide_number_tag(s, n)

# ── Slide – Chirp HC Stats ─────────────────────────────────────────────────────
def slide_chirp_hc_stats(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Chirp: HC Feature Statistics",
                       "Distribution of 13 hand-crafted features across WT vs 5xFAD groups")
    add_image(s, cf('10_chirp_hc_stats_distribution.png'),
              Inches(0.2), bar_h+Inches(0.1),
              W-Inches(0.4), H-bar_h-Inches(0.35))
    slide_number_tag(s, n)

# ── Slide – Chirp Complexity Stats ────────────────────────────────────────────
def slide_chirp_comp_stats(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Chirp: Complexity Feature Statistics",
                       "Multi-scale entropy features — group differences across scale bands")
    add_image(s, cf('11_chirp_comp_stats_distribution.png'),
              Inches(0.2), bar_h+Inches(0.1),
              W-Inches(0.4), H-bar_h-Inches(0.35))
    slide_number_tag(s, n)

# ── Slide – 4-class CNN architecture ──────────────────────────────────────────
def slide_4class_arch(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "First Approach: 4-Class CNN",
                       "BinaryCNN_NoAge backbone repurposed for 4-class problem — WT/5xFAD × Young/Adult")
    add_image(s, cf('02_4class_summary.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(8.8), H-bar_h-Inches(0.35))
    X = Inches(9.1); Y = bar_h+Inches(0.15)
    add_text(s, "Architecture", X, Y, Inches(4.0), Inches(0.38),
             size=Pt(15), bold=True, color=C_NAVY)
    Y += Inches(0.42)
    arch = [
        "Input: raw 1D ERG trace",
        "Conv1d(1→16, k=15, s=2) + BN + ReLU + MaxPool",
        "Conv1d(16→32, k=11, s=2) + BN + ReLU + MaxPool",
        "Conv1d(32→64, k=7, s=2) + BN + ReLU",
        "AdaptiveAvgPool1d(1)   ← later replaced",
        "FC: 64 → 32 → 4 classes",
        "Dropout 0.3, Adam, 80 epochs",
    ]
    for a in arch:
        add_text(s, f"  {a}", X, Y, Inches(4.0), Inches(0.4),
                 size=Pt(11.5), color=C_DARKGR)
        Y += Inches(0.4)
    add_text(s, "Key limitation: 4-class problem\nis harder with N=46; age\nconfounds genotype signal.",
             X, Y+Inches(0.1), Inches(4.0), Inches(0.9),
             size=Pt(11), color=C_MIDGR, italic=True)
    slide_number_tag(s, n)

# ── Slide – Sanity check 1 ─────────────────────────────────────────────────────
def slide_sanity1(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Sanity Check 1: Mixed-Trial Shuffle",
                       "Verifying the model learns signal structure, not trial identity")
    add_image(s, cf('02_4class_summary.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(7.5), H-bar_h-Inches(0.35))
    X = Inches(7.9); Y = bar_h+Inches(0.2)
    add_text(s, "Shuffle Test Design", X, Y, Inches(5.2), Inches(0.38),
             size=Pt(15), bold=True, color=C_NAVY)
    Y += Inches(0.46)
    pts = [
        "Randomly mix trials across subjects before\n  subject-disjoint CV split",
        "If model is memorising per-subject patterns,\n  mixed accuracy ≈ held-out accuracy",
        "If model learns general class features,\n  mixed accuracy should INCREASE",
        "Result: accuracy improves with mixed trials\n  → model generalises, not memorises",
        "Confirms the CV protocol is honest",
    ]
    for p in pts:
        add_text(s, f"• {p}", X, Y, Inches(5.2), Inches(0.72),
                 size=Pt(12), color=C_DARKGR)
        Y += Inches(0.72)
    slide_number_tag(s, n)

# ── Slide – Sanity check 2 ─────────────────────────────────────────────────────
def slide_sanity2(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Sanity Check 2: Genotype-Stratified Age Classification",
                       "Can the CNN classify Young vs Adult within each genotype separately?")
    add_image(s, cf('09_age_classification_plot.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(8.5), H-bar_h-Inches(0.35))
    X = Inches(8.8); Y = bar_h+Inches(0.2)
    add_text(s, "Rationale", X, Y, Inches(4.3), Inches(0.38),
             size=Pt(14), bold=True, color=C_NAVY)
    Y += Inches(0.45)
    pts = [
        "Age is a major source of variance in ERG",
        "A model that can't distinguish Young vs Adult\n  within ONE genotype is poorly calibrated",
        "We test CNN on WT-only and 5xFAD-only subsets",
        "Strong age classification = model captures\n  true physiological differences",
        "Acts as positive control for the CNN approach",
    ]
    for p in pts:
        add_text(s, f"• {p}", X, Y, Inches(4.3), Inches(0.7),
                 size=Pt(12), color=C_DARKGR)
        Y += Inches(0.7)
    slide_number_tag(s, n)

# ── Slide – 4-class segment results ───────────────────────────────────────────
def slide_4class_segments(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "4-Class Classification — Segment Comparison",
                       "Subject-level cross-validation accuracy across chirp segments")
    add_image(s, cf('04_chirp_segment_comparison.png'),
              Inches(0.2), bar_h+Inches(0.1),
              W-Inches(0.4), H-bar_h-Inches(0.35))
    slide_number_tag(s, n)

# ── Slide – Binary CNN architecture (ImprovedBinaryCNN) ───────────────────────
def slide_binary_arch(prs, n):
    """Draw detailed ImprovedBinaryCNN architecture diagram."""
    s = blank_slide(prs)
    bar_h = header_bar(s, "Binary CNN Architecture — ImprovedBinaryCNN  ★ Default Model",
                       "WT vs 5xFAD · Amplitude segment · Pooled AUC = 0.601")

    # ── Draw architecture flow left→right ─────────────────────────────────────
    Y0   = bar_h + Inches(0.22)
    CH   = H - Y0 - Inches(0.35)
    MID  = Y0 + CH/2
    DIAG_TOP = Y0 + CH*0.08
    DIAG_H   = CH * 0.84

    layers = [
        # (label, sublabel, color, rel_width)
        ("Input\n[B,1,2750]",          "Amplitude\nsegment",         RGBColor(0x34,0x49,0x5E), 0.7),
        ("Instance\nNorm1d",           "Per-sample\nz-score",        RGBColor(0x6C,0x3A,0x83), 1.0),
        ("Conv1d\n1→8, k=15, s=2",     "+BN+ReLU\n+MaxPool4\n→[8,343]", C_NAVY,              1.3),
        ("Conv1d\n8→16, k=11, s=2",    "+BN+ReLU\n+MaxPool4\n→[16,43]", C_NAVY,              1.3),
        ("Conv1d\n16→32, k=7, s=2",    "+BN+ReLU\n→[32,22]",        C_NAVY,                  1.3),
        ("Temporal\nStatPool",         "mean·max·std\n→[96]",        RGBColor(0x16,0x7D,0x3F),1.1),
        ("Linear\n96→32",              "+ReLU\n+Drop(0.5)",          RGBColor(0xBF,0x5C,0x00),1.0),
        ("Linear\n32→2",               "logits",                     RGBColor(0xBF,0x5C,0x00),0.8),
        ("Softmax",                    "P(WT)  P(5xFAD)",            RGBColor(0x1A,0x6E,0x30),0.9),
    ]

    total_w_units = sum(r for (_,_,_,r) in layers)
    margin_l = Inches(0.35)
    avail_w  = Inches(11.0)
    unit_w   = avail_w / total_w_units
    gap      = Inches(0.12)

    bx = margin_l
    layer_xs = []
    for (lbl, sub, col, rw) in layers:
        bw = unit_w * rw
        bh = DIAG_H * 0.8
        by = MID - bh/2
        # Box
        add_rect(s, bx, by, bw-gap, bh, fill=col)
        # Label
        add_text(s, lbl, bx+Inches(0.04), by+Inches(0.06),
                 bw-gap-Inches(0.08), bh*0.45,
                 size=Pt(10.5), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s, sub, bx+Inches(0.04), by+bh*0.5,
                 bw-gap-Inches(0.08), bh*0.44,
                 size=Pt(9.5), color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
        layer_xs.append((bx, bw-gap, by, bh))
        bx += bw

    # Arrows between blocks
    for i in range(len(layer_xs)-1):
        ax = layer_xs[i][0] + layer_xs[i][1]
        ay_mid = MID
        add_rect(s, ax, ay_mid-Inches(0.03), gap, Inches(0.06),
                 fill=C_GOLD)

    # Feature dimension annotations below boxes
    dims = ["2750","2750","343","43","22","96","32","2","2"]
    bx = margin_l
    for i, ((lx,lw,ly,lh), dim) in enumerate(zip(layer_xs, dims)):
        add_text(s, dim, lx, ly+lh+Inches(0.06), lw, Inches(0.28),
                 size=Pt(9), color=C_GOLD, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, "Feature dimension →",
             margin_l, MID+DIAG_H/2*0.8+Inches(0.38), Inches(4), Inches(0.25),
             size=Pt(10), color=C_MIDGR, italic=True)

    # Legend boxes bottom right
    Y_leg = H - Inches(0.85)
    legend_items = [
        (RGBColor(0x34,0x49,0x5E), "Input"),
        (RGBColor(0x6C,0x3A,0x83), "Normalisation"),
        (C_NAVY,                   "Conv block"),
        (RGBColor(0x16,0x7D,0x3F), "Temporal pool"),
        (RGBColor(0xBF,0x5C,0x00), "FC / Classifier"),
    ]
    xlg = W - Inches(5.5)
    for (col, lbl) in legend_items:
        add_rect(s, xlg, Y_leg+Inches(0.05), Inches(0.25), Inches(0.22), fill=col)
        add_text(s, lbl, xlg+Inches(0.3), Y_leg, Inches(0.9), Inches(0.32),
                 size=Pt(10), color=C_DARKGR)
        xlg += Inches(1.05)

    # Key improvement callout
    add_rect(s, Inches(0.25), H-Inches(0.9), Inches(5.0), Inches(0.58),
             fill=RGBColor(0xFF,0xF5,0xD6), line=C_GOLD, line_w=Pt(1.5))
    add_text(s, "★  Key fix: TemporalStatPool [mean, max, std] replaces AdaptiveAvgPool1d(1).\n"
               "    InstanceNorm1d removes per-preparation amplitude drift.  Channels: 8/16/32.",
             Inches(0.35), H-Inches(0.88), Inches(4.8), Inches(0.54),
             size=Pt(10), color=C_DARKGR)
    slide_number_tag(s, n)

# ── Slide – Age effect ─────────────────────────────────────────────────────────
def slide_age_effect(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Age Effect on CNN Classification",
                       "Does adding age metadata improve the binary WT vs 5xFAD classifier?")
    add_image(s, cf('03_age_effect_comparison.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(8.5), H-bar_h-Inches(0.35))
    X = Inches(8.8); Y = bar_h+Inches(0.2)
    add_text(s, "Age Conditions Tested", X, Y, Inches(4.3), Inches(0.38),
             size=Pt(14), bold=True, color=C_NAVY)
    Y += Inches(0.46)
    rows = [
        ("No age",         "0.601", "Baseline — no age info"),
        ("Binary age",     "0.658", "Young/Adult flag concatenated"),
        ("Continuous age", "0.692", "Age in days, normalised"),
    ]
    for (cond, auc, note) in rows:
        add_rect(s, X, Y, Inches(4.3), Inches(0.8),
                 fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(0.8))
        add_text(s, cond, X+Inches(0.1), Y+Inches(0.04), Inches(1.8), Inches(0.35),
                 size=Pt(12.5), bold=True, color=C_NAVY)
        add_text(s, f"AUC {auc}", X+Inches(2.0), Y+Inches(0.04), Inches(1.5), Inches(0.35),
                 size=Pt(13), bold=True, color=C_GOLD)
        add_text(s, note, X+Inches(0.1), Y+Inches(0.44), Inches(4.0), Inches(0.3),
                 size=Pt(11), color=C_MIDGR)
        Y += Inches(0.88)
    add_text(s, "★  Age adds modest signal.\n"
               "   Default model uses no age\n"
               "   (more generalisable).",
             X, Y+Inches(0.1), Inches(4.3), Inches(0.7),
             size=Pt(11.5), color=C_DARKGR)
    slide_number_tag(s, n)

# ── Slide – Training curves ────────────────────────────────────────────────────
def slide_training(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Binary CNN: Training Dynamics",
                       "Loss and AUC across 80 epochs for 5-fold cross-validation")
    add_image(s, cf('05_cnn_training_curves.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(8.8), H-bar_h-Inches(0.35))
    X = Inches(9.1); Y = bar_h+Inches(0.2)
    add_text(s, "Training Setup", X, Y, Inches(4.0), Inches(0.38),
             size=Pt(14), bold=True, color=C_NAVY)
    Y += Inches(0.44)
    items = [
        "Optimiser: Adam (lr=5e-4, wd=1e-2)",
        "Scheduler: CosineAnnealingLR (T=80)",
        "Loss: CrossEntropy + label_smooth=0.05",
        "Batch size: 16  (drop_last=True)",
        "Augmentation: Gaussian noise ×0.01\n  + random sign flip",
        "80 epochs, 5-fold subject-disjoint CV",
        "Best epoch selected by val. AUC",
        "Pooled AUC reported (honest metric)",
    ]
    for item in items:
        add_text(s, f"• {item}", X, Y, Inches(4.0), Inches(0.54),
                 size=Pt(11.5), color=C_DARKGR)
        Y += Inches(0.54)
    slide_number_tag(s, n)

# ── Slide – Segment comparison ─────────────────────────────────────────────────
def slide_segment_comparison(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Binary CNN: Segment Comparison",
                       "Original (BinaryCNN_NoAge) vs Improved (ImprovedBinaryCNN) across segments")
    add_image(s, cf('12_improved_chirp_cnn.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(9.0), H-bar_h-Inches(0.35))
    X = Inches(9.3); Y = bar_h+Inches(0.18)
    add_text(s, "Pooled AUC Summary", X, Y, Inches(3.8), Inches(0.38),
             size=Pt(14), bold=True, color=C_NAVY)
    Y += Inches(0.44)
    seg_data = [
        ("amplitude",      "0.601 ↑", C_GOLD),
        ("flash",          "0.567 ↑", C_GREEN),
        ("amplitude_norm", "0.546",   C_MIDGR),
        ("frequency",      "0.531",   C_MIDGR),
        ("full",           "0.520",   C_MIDGR),
    ]
    for (seg, auc, col) in seg_data:
        add_rect(s, X, Y, Inches(3.8), Inches(0.62),
                 fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(0.6))
        add_text(s, seg, X+Inches(0.1), Y+Inches(0.08), Inches(2.2), Inches(0.46),
                 size=Pt(12), color=C_NAVY)
        add_text(s, auc, X+Inches(2.3), Y+Inches(0.08), Inches(1.3), Inches(0.46),
                 size=Pt(13), bold=True, color=col)
        Y += Inches(0.66)
    add_text(s,
             "Amplitude segment is best:\n"
             "1 Hz amplitude ramp — CNN\n"
             "detects muted gain tracking.",
             X, Y+Inches(0.1), Inches(3.8), Inches(0.8),
             size=Pt(11), color=C_DARKGR, italic=True)
    slide_number_tag(s, n)

# ── Slide – HC ML ──────────────────────────────────────────────────────────────
def slide_hc_ml(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Chirp: Hand-Crafted ML Performance",
                       "Logistic Regression on 13 physiological features — 5-fold subject-disjoint CV")
    add_image(s, cf('06_handcrafted_comparison.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(7.0), H-bar_h-Inches(0.35))
    add_image(s, cf('06_handcrafted_model_selection.png'),
              Inches(7.2), bar_h+Inches(0.1),
              Inches(5.9), H-bar_h-Inches(1.3))
    # Metrics badges
    Y = H - Inches(1.1)
    for (lbl, val) in [("AUC","0.735"),("Acc","69.6%"),("Sens","60.9%"),("Spec","78.3%")]:
        metric_badge(s, lbl, val, Inches(7.2)+(["AUC","Acc","Sens","Spec"].index(lbl))*Inches(1.5),
                     Y, Inches(1.4), Inches(0.85))
    slide_number_tag(s, n)

# ── Slide – Complexity ML ──────────────────────────────────────────────────────
def slide_comp_ml(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Chirp: Complexity ML Performance",
                       "Sample Entropy (MSE) features — logistic regression / SVM")
    add_image(s, cf('07_complexity_comparison.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(7.0), H-bar_h-Inches(0.35))
    add_image(s, cf('07_complexity_model_selection.png'),
              Inches(7.2), bar_h+Inches(0.1),
              Inches(5.9), H-bar_h-Inches(1.3))
    Y = H - Inches(1.1)
    for (lbl, val) in [("AUC","0.701"),("Acc","63.0%"),("Sens","56.5%"),("Spec","69.6%")]:
        metric_badge(s, lbl, val, Inches(7.2)+(["AUC","Acc","Sens","Spec"].index(lbl))*Inches(1.5),
                     Y, Inches(1.4), Inches(0.85))
    slide_number_tag(s, n)

# ── Slide – Chirp all methods comparison ──────────────────────────────────────
def slide_chirp_comparison(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Chirp: All Approaches Compared",
                       "CNN vs Handcrafted vs Complexity — added value of each approach")
    add_image(s, cf('08_chirp_comparison.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(8.5), H-bar_h-Inches(0.35))
    X = Inches(8.8); Y = bar_h+Inches(0.18)
    add_text(s, "Summary", X, Y, Inches(4.3), Inches(0.38),
             size=Pt(15), bold=True, color=C_NAVY)
    Y += Inches(0.44)
    data = [
        ("HC + Gain Tracking",   "0.810", C_GOLD,   "★ Best chirp result"),
        ("Handcrafted (LR)",     "0.735", C_GREEN,  "Robust, interpretable"),
        ("ImprovedCNN (ampl.)",  "0.601", C_BLUE,   "End-to-end, amplitude"),
        ("Complexity (SVM)",     "0.701", C_BLUE,   "Nonlinear dynamics"),
        ("Original CNN",         "0.529", C_MIDGR,  "Near chance — deprecated"),
    ]
    for (method, auc, col, note) in data:
        add_rect(s, X, Y, Inches(4.3), Inches(0.76),
                 fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(0.6))
        add_text(s, method, X+Inches(0.1), Y+Inches(0.05), Inches(2.5), Inches(0.35),
                 size=Pt(11.5), bold=True, color=C_NAVY)
        add_text(s, f"AUC {auc}", X+Inches(2.6), Y+Inches(0.05), Inches(1.5), Inches(0.35),
                 size=Pt(13), bold=True, color=col)
        add_text(s, note, X+Inches(0.1), Y+Inches(0.44), Inches(4.0), Inches(0.26),
                 size=Pt(10.5), color=C_MIDGR, italic=True)
        Y += Inches(0.8)
    slide_number_tag(s, n)

# ── NI slides ──────────────────────────────────────────────────────────────────
def slide_ni_performance(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Natural Image: CNN Performance",
                       "ImprovedNICNN_NoAge — 10-trial average per subject · 5-fold CV")
    add_image(s, nf('01_ni_cnn_performance_bar.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(6.8), H-bar_h-Inches(0.35))
    add_image(s, nf('01_ni_cnn_training_curves.png'),
              Inches(7.0), bar_h+Inches(0.1),
              Inches(6.1), H-bar_h-Inches(0.35))
    slide_number_tag(s, n)

def slide_ni_improved(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Natural Image: Original vs Improved CNN",
                       "TemporalStatPool + InstanceNorm gives massive +0.218 AUC improvement")
    add_image(s, nf('07_improved_ni_cnn.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(9.0), H-bar_h-Inches(0.35))
    X = Inches(9.3); Y = bar_h+Inches(0.3)
    for (lbl, orig, impr) in [
        ("Pooled AUC", "0.550", "0.768"),
        ("Fold-mean AUC", "0.660", "0.770"),
        ("Pooled Acc", "~56%", "~73%"),
    ]:
        add_text(s, lbl, X, Y, Inches(3.8), Inches(0.33),
                 size=Pt(12), bold=True, color=C_NAVY)
        add_text(s, f"Original:  {orig}", X+Inches(0.1), Y+Inches(0.35),
                 Inches(1.8), Inches(0.3), size=Pt(11.5), color=C_MIDGR)
        add_text(s, f"Improved: {impr}", X+Inches(2.0), Y+Inches(0.35),
                 Inches(1.8), Inches(0.3), size=Pt(11.5), bold=True, color=C_GOLD)
        Y += Inches(0.85)
    add_text(s, "+0.218 AUC gain\nfrom architecture fixes",
             X+Inches(0.4), Y+Inches(0.2), Inches(3.0), Inches(0.7),
             size=Pt(18), bold=True, color=C_GREEN)
    slide_number_tag(s, n)

def slide_ni_hc(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Natural Image: Hand-Crafted ML",
                       "Logistic Regression on NI features — HC leads CNN in AUC")
    add_image(s, nf('02_ni_handcrafted_comparison.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(7.0), H-bar_h-Inches(0.35))
    add_image(s, nf('02_ni_handcrafted_model_selection.png'),
              Inches(7.2), bar_h+Inches(0.1),
              Inches(5.9), H-bar_h-Inches(1.3))
    Y = H - Inches(1.1)
    for (lbl, val) in [("AUC","0.860"),("Acc","70.8%"),("Sens","53.8%"),("Spec","90.9%")]:
        metric_badge(s, lbl, val, Inches(7.2)+(["AUC","Acc","Sens","Spec"].index(lbl))*Inches(1.5),
                     Y, Inches(1.4), Inches(0.85))
    slide_number_tag(s, n)

def slide_ni_comp(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Natural Image: Complexity ML",
                       "Multi-scale entropy features on NI signal")
    add_image(s, nf('03_ni_complexity_comparison.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(7.0), H-bar_h-Inches(0.35))
    add_image(s, nf('03_ni_complexity_model_selection.png'),
              Inches(7.2), bar_h+Inches(0.1),
              Inches(5.9), H-bar_h-Inches(1.3))
    Y = H - Inches(1.1)
    for (lbl, val) in [("AUC","0.829"),("Acc","70.8%"),("Sens","57.7%"),("Spec","86.4%")]:
        metric_badge(s, lbl, val, Inches(7.2)+(["AUC","Acc","Sens","Spec"].index(lbl))*Inches(1.5),
                     Y, Inches(1.4), Inches(0.85))
    slide_number_tag(s, n)

def slide_ni_comparison(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Natural Image: All Approaches Compared",
                       "HC leads, CNN competitive — complexity outperforms simple baselines")
    add_image(s, nf('04_ni_comparison.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(8.5), H-bar_h-Inches(0.35))
    X = Inches(8.8); Y = bar_h+Inches(0.3)
    add_text(s, "NI Summary", X, Y, Inches(4.3), Inches(0.38),
             size=Pt(15), bold=True, color=C_NAVY)
    Y += Inches(0.44)
    data = [
        ("HC (LR)",             "0.860", C_GOLD,  "★ Best — high specificity"),
        ("Complexity (LR)",     "0.829", C_GREEN, "Strong MSE signal"),
        ("Improved NI CNN",     "0.768", C_BLUE,  "+0.218 vs original"),
        ("Original NI CNN",     "0.550", C_MIDGR, "Near chance — fixed"),
    ]
    for (m, a, c, note) in data:
        add_rect(s, X, Y, Inches(4.3), Inches(0.76),
                 fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(0.6))
        add_text(s, m, X+Inches(0.1), Y+Inches(0.05), Inches(2.4), Inches(0.35),
                 size=Pt(11.5), bold=True, color=C_NAVY)
        add_text(s, f"AUC {a}", X+Inches(2.5), Y+Inches(0.05), Inches(1.6), Inches(0.35),
                 size=Pt(13), bold=True, color=c)
        add_text(s, note, X+Inches(0.1), Y+Inches(0.44), Inches(4.0), Inches(0.26),
                 size=Pt(10.5), color=C_MIDGR, italic=True)
        Y += Inches(0.8)
    slide_number_tag(s, n)

# ── Dual / Fusion slides ────────────────────────────────────────────────────────
def slide_dual_arch(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Combined Dual-Input Architecture",
                       "Separate CNN encoders for chirp + NI streams, fused at FC layer")
    add_image(s, mf('Fig5_ML_Classification','Figure5_ML_Classification.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(9.5), H-bar_h-Inches(0.35))
    X = Inches(9.8); Y = bar_h+Inches(0.3)
    add_text(s, "Architecture", X, Y, Inches(3.3), Inches(0.38),
             size=Pt(14), bold=True, color=C_NAVY)
    Y += Inches(0.44)
    pts = [
        "Chirp branch: ImprovedBinaryCNN\n  backbone (amplitude segment)",
        "NI branch: ImprovedNICNN_NoAge\n  backbone (trial average)",
        "Feature vectors concatenated\n  before final FC classifier",
        "Trained jointly end-to-end",
        "Same 5-fold subject-disjoint CV",
    ]
    for p in pts:
        add_text(s, f"• {p}", X, Y, Inches(3.3), Inches(0.68),
                 size=Pt(11.5), color=C_DARKGR)
        Y += Inches(0.68)
    slide_number_tag(s, n)

def slide_fusion_results(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Dual Fusion: Is Multi-Modal Better?",
                       "Combining chirp + NI in a joint CNN — rigorous 5-fold evaluation")
    add_image(s, nf('ROC_Value_of_ML.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(8.0), H-bar_h-Inches(0.35))
    X = Inches(8.3); Y = bar_h+Inches(0.2)
    add_text(s, "Multi-Modal Results", X, Y, Inches(4.8), Inches(0.38),
             size=Pt(14), bold=True, color=C_NAVY)
    Y += Inches(0.46)
    data = [
        ("Dual CNN (Fusion)", "Both",  "0.756", "0.700"),
        ("Improved NI CNN",  "NI",    "0.733", "0.768"),
        ("Improved Chirp CNN","Chirp","0.630", "0.601"),
    ]
    hdrs = ["Method","Stim.","Acc","AUC"]
    col_w = [Inches(2.0), Inches(0.7), Inches(0.7), Inches(0.7)]
    cx = X
    for h, cw in zip(hdrs, col_w):
        add_rect(s, cx, Y, cw-Inches(0.03), Inches(0.38), fill=C_NAVY)
        add_text(s, h, cx+Inches(0.05), Y+Inches(0.04), cw-Inches(0.1), Inches(0.3),
                 size=Pt(11), bold=True, color=C_WHITE)
        cx += cw
    Y += Inches(0.38)
    for i, row in enumerate(data):
        bg = C_OFFWT if i%2==0 else C_WHITE
        if row[2] == "0.756": bg = RGBColor(0xFF,0xF5,0xD6)
        cx = X
        for cell, cw in zip(row, col_w):
            add_rect(s, cx, Y, cw-Inches(0.03), Inches(0.5),
                     fill=bg, line=C_LIGHTGR, line_w=Pt(0.5))
            add_text(s, cell, cx+Inches(0.05), Y+Inches(0.08), cw-Inches(0.1), Inches(0.34),
                     size=Pt(11), color=C_DARKGR)
            cx += cw
        Y += Inches(0.5)
    Y += Inches(0.2)
    add_text(s,
             "★  Dual fusion improves accuracy (+0.023)\n"
             "    but NI CNN alone leads in AUC (0.768).\n"
             "    Fusion provides complementary information,\n"
             "    reduces variance (Std 0.083 vs 0.113).",
             X, Y, Inches(4.8), Inches(1.2),
             size=Pt(12), color=C_DARKGR)
    slide_number_tag(s, n)

# ── Interpretability slides ─────────────────────────────────────────────────────
def slide_interp_overview(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "What Can We Learn From the CNN?",
                       "Four complementary interpretability approaches applied to ImprovedBinaryCNN")
    approaches = [
        ("Learned Kernels", cf('13_a_kernels.png'),
         "Conv1 kernels visualised — bandpass / derivative shapes"),
        ("GradCAM / Integrated Gradients", cf('13_b_gradcam.png'),
         "Temporal importance maps — which time points drive decisions?"),
        ("Temporal Attention Maps", cf('14_b_attention_maps_amplitude.png'),
         "AttentionBinaryCNN weights — direct interpretable heatmap"),
        ("Optimal Stimuli & PCA", cf('15_a_optimal_signals.png'),
         "Gradient ascent reveals what the CNN considers 'ideal 5xFAD'"),
    ]
    xpos = [Inches(0.2), Inches(3.55), Inches(6.9), Inches(10.25)]
    bw   = Inches(3.2)
    for (title, fig, desc), xp in zip(approaches, xpos):
        add_image(s, fig, xp, bar_h+Inches(0.12), bw, Inches(4.2))
        add_rect(s, xp, bar_h+Inches(4.35), bw, Inches(0.9), fill=C_NAVY)
        add_text(s, title, xp+Inches(0.08), bar_h+Inches(4.38), bw-Inches(0.16), Inches(0.4),
                 size=Pt(11), bold=True, color=C_GOLD)
        add_text(s, desc, xp+Inches(0.08), bar_h+Inches(4.78), bw-Inches(0.16), Inches(0.36),
                 size=Pt(9.5), color=C_WHITE)
    slide_number_tag(s, n)

def slide_kernels_gradients(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Learned Representations: Kernels & Gradient Methods",
                       "First-layer kernels + GradCAM + Integrated Gradients on amplitude segment")
    add_image(s, cf('13_a_kernels.png'),
              Inches(0.2), bar_h+Inches(0.1), Inches(4.2), H-bar_h-Inches(0.35))
    add_image(s, cf('13_b_gradcam.png'),
              Inches(4.5), bar_h+Inches(0.1), Inches(4.5), H-bar_h-Inches(0.35))
    add_image(s, cf('13_c_integrated_grads.png'),
              Inches(9.1), bar_h+Inches(0.1), Inches(4.0), H-bar_h-Inches(0.35))
    slide_number_tag(s, n)

def slide_attention_maps(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Temporal Attention Maps — AttentionBinaryCNN",
                       "Learned attention weights localise which time points distinguish WT from 5xFAD")
    add_image(s, cf('14_b_attention_maps_amplitude.png'),
              Inches(0.2), bar_h+Inches(0.1), Inches(6.5), H-bar_h-Inches(0.35))
    add_image(s, cf('14_c_attention_stats_amplitude.png'),
              Inches(6.8), bar_h+Inches(0.1), Inches(6.3), H-bar_h-Inches(0.35))
    slide_number_tag(s, n)

def slide_optimal_pca(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Optimal Signals & PCA Landscape",
                       "Gradient ascent reveals the CNN's internal representation of each class")
    add_image(s, cf('15_a_optimal_signals.png'),
              Inches(0.2), bar_h+Inches(0.1), Inches(6.5), H-bar_h-Inches(0.35))
    add_image(s, cf('15_b_pca_coefficients.png'),
              Inches(6.8), bar_h+Inches(0.1), Inches(6.3), H-bar_h-Inches(0.35))
    slide_number_tag(s, n)

def slide_phase_envelope(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Phase & Envelope Analysis",
                       "Fundamental-frequency phase and amplitude envelope drive discrimination")
    add_image(s, cf('16_a_phase_motivation.png'),
              Inches(0.2), bar_h+Inches(0.1), Inches(4.8), H-bar_h-Inches(0.35))
    add_image(s, cf('16_b_feature_distributions.png'),
              Inches(5.1), bar_h+Inches(0.1), Inches(4.0), H-bar_h-Inches(0.35))
    add_image(s, cf('16_c_classification.png'),
              Inches(9.2), bar_h+Inches(0.1), Inches(3.9), H-bar_h-Inches(0.35))
    slide_number_tag(s, n)

# ── Muted gain tracking hypothesis ────────────────────────────────────────────
def slide_gain_hypothesis(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Muted Gain Tracking: Biological Hypothesis",
                       "5xFAD retinas show reduced gain at the fundamental frequency of the amplitude chirp")
    Y = bar_h + Inches(0.2)

    # Left: explain what gain tracking is
    add_rect(s, Inches(0.25), Y, Inches(6.3), H-Y-Inches(0.35),
             fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(1))
    add_text(s, "What is Gain Tracking?", Inches(0.4), Y+Inches(0.1),
             Inches(6.0), Inches(0.4), size=Pt(15), bold=True, color=C_NAVY)
    items = [
        "The amplitude chirp segment (24–35 s) is a\n"
        "  1 Hz sine wave with LINEARLY INCREASING amplitude",
        "Healthy retinas track the rising amplitude gain:\n"
        "  their ERG response envelope grows proportionally",
        "5xFAD retinas show REDUCED gain tracking:\n"
        "  the response envelope is 'muted' — less growth",
        "This reflects impaired outer/middle retinal pathway\n"
        "  processing of temporally modulated light stimuli",
        "CNN picks up this pattern without being told explicitly",
    ]
    yy = Y + Inches(0.6)
    for item in items:
        add_text(s, f"• {item}", Inches(0.4), yy, Inches(5.9), Inches(0.72),
                 size=Pt(12.5), color=C_DARKGR)
        yy += Inches(0.72)
    add_image(s, cf('16_d_envelope.png'),
              Inches(0.3), yy, Inches(6.1), H-yy-Inches(0.4))

    # Right: evidence
    add_rect(s, Inches(6.7), Y, Inches(6.38), H-Y-Inches(0.35),
             fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(1))
    add_text(s, "Three Levels of Evidence", Inches(6.85), Y+Inches(0.1),
             Inches(6.0), Inches(0.4), size=Pt(15), bold=True, color=C_NAVY)
    levels = [
        ("Level 1 — Direct Biology",
         "Hilbert envelope of bandpass-filtered (0.5–1.5 Hz)\n"
         "amplitude signal compared between groups.\n"
         "Result: 0/50 time-bins FDR-significant (N=46 underpowered).",
         C_MIDGR),
        ("Level 2 — Corrected Virtual Blockade (CNN)",
         "Remove fundamental band → measure asymmetric drop in P(5xFAD).\n"
         "5xFAD loses 2× more classification probability than WT.\n"
         "→ Fundamental is group-specifically encoded in the CNN.",
         C_GOLD),
        ("Level 3 — HC Classifier Improvement",
         "Add 4 gain-tracking features to HC feature set.\n"
         "AUC: 0.735 → 0.810  (+0.075, +10.2% relative).\n"
         "→ Gain pattern is independently predictive.",
         C_GREEN),
    ]
    yy = Y + Inches(0.6)
    for (title, body, col) in levels:
        add_rect(s, Inches(6.75), yy, Inches(0.1), Inches(1.4), fill=col)
        add_text(s, title, Inches(6.95), yy+Inches(0.04), Inches(6.0), Inches(0.38),
                 size=Pt(12.5), bold=True, color=col)
        add_text(s, body, Inches(6.95), yy+Inches(0.42), Inches(6.0), Inches(0.9),
                 size=Pt(11.5), color=C_DARKGR)
        yy += Inches(1.55)
    slide_number_tag(s, n)

# ── Corrected virtual blockade – math ─────────────────────────────────────────
def slide_blockade_math(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Level 2: Corrected Virtual Blockade — Full Mathematical Proof",
                       "Why removing the 1 Hz band proves group-specific encoding by the CNN")
    Y = bar_h + Inches(0.18)

    # Top: why naive blockade is confounded
    add_rect(s, Inches(0.2), Y, Inches(13.0), Inches(1.1),
             fill=RGBColor(0xFF,0xF0,0xF0), line=C_RED, line_w=Pt(1.2))
    add_text(s, "⚠  Why the naive virtual blockade is insufficient:",
             Inches(0.35), Y+Inches(0.06), Inches(8), Inches(0.35),
             size=Pt(12.5), bold=True, color=C_RED)
    add_text(s,
             "The amplitude chirp IS a 1 Hz sinusoid. Removing 1 Hz from a 1 Hz response will always hurt performance — for BOTH groups — because "
             "it removes the primary stimulus-driven component. This is trivially expected and proves nothing about group differences.",
             Inches(0.35), Y+Inches(0.45), Inches(12.6), Inches(0.55),
             size=Pt(11.5), color=C_DARKGR)
    Y += Inches(1.25)

    # Left column: math
    CX = Inches(0.2); CW = Inches(7.9)
    add_rect(s, CX, Y, CW, H-Y-Inches(0.35),
             fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(1))
    add_text(s, "Mathematical Framework", CX+Inches(0.15), Y+Inches(0.1),
             CW-Inches(0.3), Inches(0.38), size=Pt(14), bold=True, color=C_NAVY)
    yy = Y + Inches(0.55)
    math_items = [
        ("Notation",
         "Let xᵢ ∈ ℝᵀ  be the mean ERG signal for subject i  (T = 2750 samples, f_s = 250 Hz)\n"
         "Let xᵢ^block  =  xᵢ  −  bandpass(xᵢ, 0.5–1.5 Hz)  [fundamental removed]\n"
         "Let f_θ(·) = ensemble of 5 ImprovedBinaryCNN models, output: P(5xFAD | ·)"),
        ("Define per-subject blockade effect",
         "ΔPᵢ  =  f_θ(xᵢ^block) − f_θ(xᵢ)\n"
         "This is signed: ΔP < 0 means removing the fundamental REDUCES P(5xFAD)."),
        ("Null hypothesis  H₀",
         "E[ΔP | WT] = E[ΔP | 5xFAD]\n"
         "i.e., the fundamental is equally informative for both groups."),
        ("Alternative hypothesis  H₁ (muted gain tracking)",
         "E[ΔP | 5xFAD] ≪ E[ΔP | WT]  (both negative, but 5xFAD effect much larger)\n"
         "i.e., the CNN relies MORE on the fundamental to recognise 5xFAD subjects."),
        ("Critical distinction from naive blockade",
         "We do not test whether 'removing 1 Hz hurts' (trivially true for both).\n"
         "We test whether it hurts ASYMMETRICALLY  →  group-specific encoding."),
    ]
    for (title, body) in math_items:
        add_rect(s, CX+Inches(0.12), yy, Inches(0.08), Inches(0.3), fill=C_NAVY)
        add_text(s, title, CX+Inches(0.3), yy, CW-Inches(0.45), Inches(0.3),
                 size=Pt(12), bold=True, color=C_NAVY)
        add_text(s, body, CX+Inches(0.3), yy+Inches(0.33), CW-Inches(0.45),
                 Inches(0.7), size=Pt(11), color=C_DARKGR)
        yy += Inches(1.05)

    # Right column: filter design + result
    RX = Inches(8.3); RW = Inches(4.85)
    add_rect(s, RX, Y, RW, H-Y-Inches(0.35),
             fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(1))
    add_text(s, "Blockade Filter Design", RX+Inches(0.15), Y+Inches(0.1),
             RW-Inches(0.3), Inches(0.38), size=Pt(14), bold=True, color=C_NAVY)
    yy = Y + Inches(0.55)
    filter_pts = [
        "Zero-phase Butterworth bandpass 0.5–1.5 Hz",
        "Applied via scipy.signal.filtfilt (no phase lag)",
        "Subtracted from signal to isolate ~1 Hz component",
        "Removes fundamental while preserving all other bands\n  (drift <0.5 Hz, harmonics >1.5 Hz, broadband noise)",
    ]
    for p in filter_pts:
        add_text(s, f"• {p}", RX+Inches(0.15), yy, RW-Inches(0.3), Inches(0.58),
                 size=Pt(11), color=C_DARKGR)
        yy += Inches(0.58)
    yy += Inches(0.15)
    add_rect(s, RX+Inches(0.1), yy, RW-Inches(0.2), Inches(2.5),
             fill=C_NAVY)
    add_text(s, "RESULT", RX+Inches(0.3), yy+Inches(0.1), RW-Inches(0.4), Inches(0.4),
             size=Pt(14), bold=True, color=C_GOLD)
    result_lines = [
        ("WT   mean ΔP =",  "−0.105"),
        ("5xFAD mean ΔP =", "−0.211"),
        ("Asymmetry ratio =","2.0×"),
    ]
    yy2 = yy + Inches(0.55)
    for (lbl, val) in result_lines:
        add_text(s, lbl, RX+Inches(0.3), yy2, RW-Inches(1.3), Inches(0.42),
                 size=Pt(13), color=C_WHITE)
        add_text(s, val, RX+Inches(0.3)+(RW-Inches(0.4))*0.55, yy2, Inches(2.0), Inches(0.42),
                 size=Pt(14), bold=True, color=C_GOLD)
        yy2 += Inches(0.55)
    add_text(s,
             "→ Fundamental encodes\n5xFAD-specific patterns\n2× more than WT.",
             RX+Inches(0.3), yy2, RW-Inches(0.4), Inches(0.65),
             size=Pt(12.5), bold=True, color=C_GREEN)
    slide_number_tag(s, n)

# ── Blockade results figure ────────────────────────────────────────────────────
def slide_blockade_results(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Level 2: Virtual Blockade Results",
                       "ImprovedBinaryCNN (AUC=0.601) — asymmetric fundamental contribution confirmed")
    add_image(s, cf('24_muted_gain_tracking.png'),
              Inches(0.2), bar_h+Inches(0.1),
              W-Inches(0.4), H-bar_h-Inches(0.35))
    slide_number_tag(s, n)

# ── Gain tracking HC improvement ──────────────────────────────────────────────
def slide_gain_hc_improvement(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Level 3: CNN Knowledge → HC Feature Engineering",
                       "Gain-tracking features motivated by CNN interpretability findings")
    add_image(s, cf('23_gain_tracking_comparison.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(8.5), H-bar_h-Inches(0.35))
    X = Inches(8.8); Y = bar_h+Inches(0.2)
    add_text(s, "4 New Gain-Tracking Features", X, Y, Inches(4.3), Inches(0.38),
             size=Pt(14), bold=True, color=C_NAVY)
    Y += Inches(0.44)
    feats = [
        ("AmpFund_RMS",        "RMS of bandpass(0.5–1.5 Hz)\namplitude segment"),
        ("AmpFund_frac",       "Fraction of total RMS in\nthe fundamental band"),
        ("AmpEnv_slope_norm",  "Normalised slope of Hilbert\nenvelope (gain trajectory)"),
        ("AmpEnv_late_early",  "Late/early envelope ratio\n(does gain grow or plateau?)"),
    ]
    for (nm, desc) in feats:
        add_rect(s, X, Y, Inches(4.3), Inches(0.75),
                 fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(0.7))
        add_text(s, nm, X+Inches(0.1), Y+Inches(0.05), Inches(4.0), Inches(0.3),
                 size=Pt(12), bold=True, color=C_NAVY)
        add_text(s, desc, X+Inches(0.1), Y+Inches(0.38), Inches(4.0), Inches(0.3),
                 size=Pt(11), color=C_MIDGR)
        Y += Inches(0.79)
    Y += Inches(0.1)
    add_rect(s, X, Y, Inches(4.3), Inches(1.0), fill=RGBColor(0xFF,0xF5,0xD6),
             line=C_GOLD, line_w=Pt(1.5))
    add_text(s, "AUC improvement",
             X+Inches(0.15), Y+Inches(0.06), Inches(4.0), Inches(0.32),
             size=Pt(13), bold=True, color=C_NAVY)
    add_text(s, "HC only          → 0.735\nHC + GainTracking → 0.810  (+0.075)",
             X+Inches(0.15), Y+Inches(0.42), Inches(4.0), Inches(0.5),
             size=Pt(13), bold=True, color=C_GOLD)
    slide_number_tag(s, n)

# ── Gemini scripts critique ─────────────────────────────────────────────────────
def slide_gemini_critique(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "CNN Interpretability: Prior Work & Critique",
                       "Scripts 19–22: virtual blockade + counterfactual edits (Gemini analysis)")
    add_image(s, cf('19_virtual_blockade.png'),
              Inches(0.2), bar_h+Inches(0.1),
              Inches(4.3), H-bar_h-Inches(0.35))
    add_image(s, cf('21_cure_spectral_analysis.png'),
              Inches(4.6), bar_h+Inches(0.1),
              Inches(4.3), H-bar_h-Inches(0.35))
    add_image(s, cf('22_symmetrical_story.png'),
              Inches(9.0), bar_h+Inches(0.1),
              Inches(4.1), H-bar_h-Inches(0.35))

    # Critique boxes
    Y = H - Inches(2.1)
    items_agree = [
        "Script 19 (virtual blockade): Valid approach,\n  correctly identifies fundamental as key band",
        "Concept of muted gain tracking is biological\n  and mechanistically grounded",
    ]
    items_disagree = [
        "Scripts 20–22 (counterfactuals): Smoothness\n  regularizer β=50 FORCES low-freq edits artificially",
        "'Spectral symmetry at 1 Hz' finding is a\n  regularization artifact, not a genuine CNN insight",
        "Used near-chance original CNN (AUC≈0.529)\n  → findings less meaningful",
    ]
    add_rect(s, Inches(0.2), Y, Inches(6.0), Inches(1.75),
             fill=RGBColor(0xE8,0xF8,0xE8), line=C_GREEN, line_w=Pt(1.2))
    add_text(s, "✓ Agreed",
             Inches(0.35), Y+Inches(0.05), Inches(5.7), Inches(0.35),
             size=Pt(13), bold=True, color=C_GREEN)
    yy = Y+Inches(0.42)
    for item in items_agree:
        add_text(s, f"• {item}", Inches(0.35), yy, Inches(5.7), Inches(0.55),
                 size=Pt(11), color=C_DARKGR)
        yy += Inches(0.58)

    add_rect(s, Inches(6.4), Y, Inches(6.73), Inches(1.75),
             fill=RGBColor(0xFF,0xED,0xED), line=C_RED, line_w=Pt(1.2))
    add_text(s, "✗ Methodological issues",
             Inches(6.55), Y+Inches(0.05), Inches(6.4), Inches(0.35),
             size=Pt(13), bold=True, color=C_RED)
    yy = Y+Inches(0.42)
    for item in items_disagree:
        add_text(s, f"• {item}", Inches(6.55), yy, Inches(6.4), Inches(0.45),
                 size=Pt(11), color=C_DARKGR)
        yy += Inches(0.5)
    slide_number_tag(s, n)

# ── Conclusions ────────────────────────────────────────────────────────────────
def slide_conclusions(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Conclusions")
    Y = bar_h + Inches(0.2)

    cols = [
        ("Chirp CNN", [
            "ImprovedBinaryCNN achieves AUC=0.601 on\n  amplitude segment (chance=0.5)",
            "TemporalStatPool + InstanceNorm critical:\n  +0.151 AUC vs original architecture",
            "Amplitude segment is most informative\n  for binary classification",
            "Age adds modest signal (0.601→0.692)",
        ], C_NAVY),
        ("Natural Image CNN", [
            "ImprovedNICNN achieves AUC=0.768\n  — much stronger than chirp CNN",
            "HC still leads (AUC=0.860) but gap\n  narrowed from original",
            "Architecture improvements universal:\n  +0.218 AUC improvement",
            "Dataset size favours simpler models",
        ], C_BLUE),
        ("Handcrafted ML", [
            "Chirp HC: AUC=0.735 (baseline)",
            "NI HC: AUC=0.860 (best single method)",
            "HC + GainTracking: AUC=0.810 ← new best chirp",
            "CNN interpretability → feature engineering\n  → measurable AUC gain (+0.075)",
        ], C_GREEN),
        ("Multi-modal & Biology", [
            "Dual fusion improves accuracy (0.756)\n  but NI CNN leads in AUC",
            "5xFAD retinas show muted gain tracking\n  at 1 Hz in amplitude chirp",
            "Fundamental band 2× more diagnostic\n  for 5xFAD than WT (virtual blockade)",
            "Deficit in outer/middle retinal pathway\n  response gain, not just amplitude",
        ], C_GOLD),
    ]
    col_w = (W - Inches(0.5)) / 4
    for i, (title, items, col) in enumerate(cols):
        cx = Inches(0.2) + i * col_w
        add_rect(s, cx, Y, col_w-Inches(0.1), Inches(0.48), fill=col)
        add_text(s, title, cx+Inches(0.1), Y+Inches(0.06),
                 col_w-Inches(0.3), Inches(0.36),
                 size=Pt(14), bold=True, color=C_WHITE)
        yy = Y + Inches(0.55)
        for item in items:
            add_text(s, f"• {item}", cx+Inches(0.1), yy, col_w-Inches(0.25),
                     Inches(0.78), size=Pt(11), color=C_DARKGR)
            yy += Inches(0.82)
    slide_number_tag(s, n)

def slide_next_steps(prs, n):
    s = blank_slide(prs)
    bar_h = header_bar(s, "Next Steps & Open Questions")
    Y = bar_h + Inches(0.25)
    items = [
        ("Expand the dataset",
         "More animals needed to reach statistical power for Level 1 (biology) comparisons.\n"
         "Target N≥80 for FDR-significant envelope trajectory differences."),
        ("In-vivo ERG validation",
         "Translate findings from in-vitro microarray recordings to standard in-vivo ERG.\n"
         "Test whether 1 Hz gain-tracking deficit is detectable non-invasively."),
        ("Refine gain-tracking features",
         "Explore frequency-resolved gain curves across the full amplitude ramp.\n"
         "Compute time-resolved instantaneous gain using short-time Fourier transform."),
        ("Cross-stimulus feature fusion",
         "CCA already shows correlated latent structure between chirp and NI.\n"
         "Explore feature-level (not model-level) fusion for interpretable multi-modal classifier."),
        ("Earlier timepoints (young mice)",
         "Current models struggle with young animals — low AUC.\n"
         "Longitudinal tracking: can we detect AD onset before behavioural symptoms?"),
        ("Histological validation",
         "Correlate ERG gain-tracking features with retinal layer thickness (OCT) and\n"
         "amyloid deposition in the retinal ganglion cell layer."),
    ]
    bh = (H-Y-Inches(0.35)) / 3
    for i, (title, body) in enumerate(items):
        row, col = divmod(i, 2)
        cx = Inches(0.25) + col * Inches(6.6)
        cy = Y + row * bh
        add_rect(s, cx, cy, Inches(6.5), bh-Inches(0.1),
                 fill=C_OFFWT, line=C_LIGHTGR, line_w=Pt(0.8))
        add_rect(s, cx, cy, Inches(0.14), bh-Inches(0.1), fill=C_GOLD)
        add_text(s, title, cx+Inches(0.24), cy+Inches(0.08),
                 Inches(6.0), Inches(0.38), size=Pt(13), bold=True, color=C_NAVY)
        add_text(s, body, cx+Inches(0.24), cy+Inches(0.5),
                 Inches(6.0), bh-Inches(0.65), size=Pt(11.5), color=C_DARKGR)
    slide_number_tag(s, n)

# ── Build all slides ───────────────────────────────────────────────────────────
def build():
    prs = new_prs()
    n   = 1

    slide_title(prs, n); n+=1                        # 1
    slide_motivation(prs, n); n+=1                   # 2
    slide_experiment(prs, n); n+=1                   # 3
    slide_demographics(prs, n); n+=1                 # 4
    slide_stimuli(prs, n); n+=1                      # 5
    slide_chirp_segments(prs, n); n+=1               # 6
    slide_hc_features(prs, n); n+=1                  # 7
    slide_complexity(prs, n); n+=1                   # 8

    # ── Part 1: Chirp ─────────────────────────────────────────────────────────
    slide_section(prs, n, "PART 1", "Chirp Stimulus Analysis",
                  "CNN · Handcrafted ML · Complexity Features"); n+=1          # 9
    slide_chirp_hc_stats(prs, n); n+=1               # 10
    slide_chirp_comp_stats(prs, n); n+=1             # 11
    slide_4class_arch(prs, n); n+=1                  # 12
    slide_sanity1(prs, n); n+=1                      # 13
    slide_sanity2(prs, n); n+=1                      # 14
    slide_4class_segments(prs, n); n+=1              # 15
    slide_binary_arch(prs, n); n+=1                  # 16
    slide_age_effect(prs, n); n+=1                   # 17
    slide_training(prs, n); n+=1                     # 18
    slide_segment_comparison(prs, n); n+=1           # 19
    slide_hc_ml(prs, n); n+=1                        # 20
    slide_comp_ml(prs, n); n+=1                      # 21
    slide_chirp_comparison(prs, n); n+=1             # 22

    # ── Part 2: Natural Image ─────────────────────────────────────────────────
    slide_section(prs, n, "PART 2", "Natural Image Analysis",
                  "Temporal & spatial dynamics of reality responses"); n+=1    # 23
    slide_ni_improved(prs, n); n+=1                  # 24
    slide_ni_performance(prs, n); n+=1               # 25
    slide_ni_hc(prs, n); n+=1                        # 26
    slide_ni_comp(prs, n); n+=1                      # 27
    slide_ni_comparison(prs, n); n+=1                # 28

    # ── Part 3: Dual ──────────────────────────────────────────────────────────
    slide_section(prs, n, "PART 3", "The Combined Approach",
                  "Fusing chirp and natural image in a dual-input CNN"); n+=1  # 29
    slide_dual_arch(prs, n); n+=1                    # 30
    slide_fusion_results(prs, n); n+=1               # 31

    # ── Part 4: CNN Interpretability ──────────────────────────────────────────
    slide_section(prs, n, "PART 4", "What Can We Learn From the CNN?",
                  "Interpretability → biological insight → feature engineering"); n+=1  # 32
    slide_interp_overview(prs, n); n+=1              # 33
    slide_kernels_gradients(prs, n); n+=1            # 34
    slide_attention_maps(prs, n); n+=1               # 35
    slide_optimal_pca(prs, n); n+=1                  # 36
    slide_phase_envelope(prs, n); n+=1               # 37
    slide_gemini_critique(prs, n); n+=1              # 38
    slide_gain_hypothesis(prs, n); n+=1              # 39
    slide_blockade_math(prs, n); n+=1                # 40
    slide_blockade_results(prs, n); n+=1             # 41
    slide_gain_hc_improvement(prs, n); n+=1          # 42

    # ── Conclusions ───────────────────────────────────────────────────────────
    slide_section(prs, n, "CONCLUSIONS", "Summary & Next Steps",
                  "Key findings and future directions"); n+=1                  # 43
    slide_conclusions(prs, n); n+=1                  # 44
    slide_next_steps(prs, n); n+=1                   # 45

    prs.save(OUT_PATH)
    print(f"\n✓ Saved {n-1}-slide presentation → {OUT_PATH}")

if __name__ == '__main__':
    build()
